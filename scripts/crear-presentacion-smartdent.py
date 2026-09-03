from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "doc" / "SmartDent_Presentacion_Avance_1.pptx"
ASSETS = ROOT / "doc" / "ppt-assets" / "tech"
IMG = ROOT / "maquetacion-html" / "img"

NAVY = "06172E"
NAVY_2 = "0B2447"
CYAN = "11D8E5"
GOLD = "D4A900"
GOLD_LIGHT = "F8DE7E"
WHITE = "FFFFFF"
INK = "10213B"
MUTED = "60708A"
LIGHT = "F4F7FC"
LINE = "D9E2EF"
GREEN = "198754"
RED = "D7263D"


def rgb(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def add_shape(slide, shape_type, x, y, w, h, fill, line=None, radius=False):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else shape_type
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    return shape


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False,
             font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
             margin=0.04, line_spacing=1.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = rgb(color)
    return box


def add_rich_text(slide, runs, x, y, w, h, size=18, align=PP_ALIGN.LEFT,
                  valign=MSO_ANCHOR.TOP, font="Aptos"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(0.02)
    frame.margin_top = frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    for text, color, bold in runs:
        run = paragraph.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def add_picture_crop(slide, path, x, y, w, h):
    path = Path(path)
    with Image.open(path) as im:
        img_ratio = im.width / im.height
    frame_ratio = w / h
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    if img_ratio > frame_ratio:
        visible = frame_ratio / img_ratio
        pic.crop_left = pic.crop_right = (1 - visible) / 2
    else:
        visible = img_ratio / frame_ratio
        pic.crop_top = pic.crop_bottom = (1 - visible) / 2
    return pic


def add_logo(slide, x=0.45, y=0.20, w=2.05):
    path = IMG / "branding" / "logo-smartdent.png"
    with Image.open(path) as im:
        ratio = im.height / im.width
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(w * ratio))


def add_header(slide, title, subtitle=None, number=None, dark=False):
    bg = NAVY if dark else WHITE
    title_color = WHITE if dark else NAVY
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, bg)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.07, CYAN)
    add_logo(slide)
    add_text(slide, title, 0.62, 0.86, 11.8, 0.52, 27, title_color, True)
    if subtitle:
        add_text(slide, subtitle, 0.64, 1.38, 11.6, 0.42, 12.5,
                 "AFC3DB" if dark else MUTED)
    if number is not None:
        add_text(slide, f"{number:02d}", 12.25, 0.26, 0.52, 0.28, 10,
                 CYAN if dark else GOLD, True, align=PP_ALIGN.RIGHT)


def add_footer(slide, number, dark=False):
    color = "8DA2BD" if dark else "8190A6"
    add_text(slide, "DESARROLLO WEB INTEGRADO · SMARTDENT", 0.64, 7.15, 6.8, 0.2,
             8, color, True)
    add_text(slide, str(number), 12.18, 7.15, 0.55, 0.2, 8, color, True,
             align=PP_ALIGN.RIGHT)


def add_bullet_list(slide, items, x, y, w, h, size=17, color=INK, bullet_color=GOLD,
                    gap=0.58):
    for index, item in enumerate(items):
        yy = y + index * gap
        add_shape(slide, MSO_SHAPE.OVAL, x, yy + 0.10, 0.12, 0.12, bullet_color)
        add_text(slide, item, x + 0.25, yy, w - 0.25, gap, size, color)


def add_card(slide, x, y, w, h, title, body, accent=GOLD, dark=False):
    fill = NAVY_2 if dark else WHITE
    title_color = WHITE if dark else NAVY
    body_color = "C1D0E4" if dark else MUTED
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill,
              "244263" if dark else LINE, True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.08, h, accent)
    add_text(slide, title, x + 0.28, y + 0.22, w - 0.5, 0.4, 16, title_color, True)
    add_text(slide, body, x + 0.28, y + 0.70, w - 0.5, h - 0.85, 11.5, body_color,
             line_spacing=1.08)


def add_code_block(slide, code, x, y, w, h):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, "07111F", "244263", True)
    add_shape(slide, MSO_SHAPE.OVAL, x + 0.22, y + 0.18, 0.11, 0.11, "FF5F57")
    add_shape(slide, MSO_SHAPE.OVAL, x + 0.40, y + 0.18, 0.11, 0.11, "FEBC2E")
    add_shape(slide, MSO_SHAPE.OVAL, x + 0.58, y + 0.18, 0.11, 0.11, "28C840")
    add_text(slide, code, x + 0.22, y + 0.50, w - 0.44, h - 0.65, 10.2,
             "D7E6F7", font="Consolas", line_spacing=0.88)


# 1. Portada
slide = prs.slides.add_slide(blank)
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, NAVY)
add_picture_crop(slide, IMG / "heroes" / "hero-nosotros.png", 7.55, 0, 5.783, 7.5)
add_shape(slide, MSO_SHAPE.RECTANGLE, 7.47, 0, 0.12, 7.5, CYAN)
add_logo(slide, 0.72, 0.58, 2.65)
slide.shapes.add_picture(str(ROOT / "doc" / "ppt-assets" / "logo-utp.png"),
                         Inches(5.90), Inches(0.58), Inches(0.92), Inches(0.34))
add_text(slide, "SISTEMA WEB INTEGRADO", 0.76, 2.02, 5.9, 0.34, 12, GOLD_LIGHT, True)
add_text(slide, "SmartDent", 0.72, 2.42, 6.15, 0.82, 38, WHITE, True)
add_text(slide, "Gestión inteligente de citas y atenciones odontológicas", 0.76, 3.33,
         5.92, 1.0, 20, "C5D7EB", False, line_spacing=1.08)
add_shape(slide, MSO_SHAPE.RECTANGLE, 0.76, 4.55, 1.25, 0.07, GOLD)
add_text(slide, "Avance 1 · Desarrollo Web Integrado", 0.76, 4.82, 5.4, 0.36,
         13, WHITE, True)
add_text(slide, "Universidad Tecnológica del Perú · 2026", 0.76, 5.25, 5.4, 0.32,
         11.5, "9FB4CC")
add_text(slide,
         "Acevedo Huarachi Kelvin Jesus · U23309803\n"
         "Añorga Pinedo Paolo Alexander · U23305864\n"
         "Calle Paredes Maykol Adan · U23242558\n"
         "Salinas Perez Joseph Sebastian · U23325202",
         0.76, 5.90, 5.7, 1.02, 8.6, "9FB4CC", line_spacing=0.92)


# 2. Tecnologías
slide = prs.slides.add_slide(blank)
add_header(slide, "Tecnologías implementadas", "Stack utilizado en el prototipo funcional", 2)
techs = [
    ("html5", "HTML5", "Estructura web"),
    ("css3", "CSS3", "Estilos globales"),
    ("tailwindcss", "Tailwind CSS", "Componentes visuales"),
    ("javascript", "JavaScript", "Interacción y API"),
    ("java", "Java 21", "Lenguaje backend"),
    ("spring", "Spring Boot", "API REST"),
    ("hibernate", "JPA / Hibernate", "Persistencia ORM"),
    ("mariadb", "MariaDB", "Base de datos"),
    ("junit", "JUnit 5", "Pruebas"),
    ("swagger", "Swagger", "Documentación API"),
]
for idx, (icon, title, desc) in enumerate(techs):
    col, row = idx % 5, idx // 5
    x, y = 0.62 + col * 2.48, 2.0 + row * 2.12
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 2.14, 1.72, WHITE, LINE, True)
    slide.shapes.add_picture(str(ASSETS / f"{icon}.png"), Inches(x + 0.73), Inches(y + 0.18),
                             Inches(0.68), Inches(0.68))
    add_text(slide, title, x + 0.12, y + 0.94, 1.90, 0.28, 12.5, NAVY, True,
             align=PP_ALIGN.CENTER)
    add_text(slide, desc, x + 0.12, y + 1.28, 1.90, 0.22, 9.5, MUTED,
             align=PP_ALIGN.CENTER)
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.62, 6.35, 12.1, 0.48, "E9FBFD", "BCEFF3", True)
add_text(slide, "Seguridad transversal: Spring Security + JWT + BCrypt", 0.88, 6.46, 11.55,
         0.22, 11.5, NAVY, True, align=PP_ALIGN.CENTER)
add_footer(slide, 2)


# 3. Problema
slide = prs.slides.add_slide(blank)
add_header(slide, "El problema que resolvemos", "La gestión manual limita el crecimiento de la clínica", 3, True)
add_picture_crop(slide, IMG / "servicios" / "consulta-diagnostico.png", 7.82, 1.88, 4.88, 4.72)
add_shape(slide, MSO_SHAPE.RECTANGLE, 7.82, 6.28, 4.88, 0.32, GOLD)
add_card(slide, 0.65, 2.00, 6.35, 1.10, "01  Cruces de horarios",
         "Reservas coordinadas por llamadas, cuadernos o mensajes generan duplicidad.", CYAN, True)
add_card(slide, 0.65, 3.33, 6.35, 1.10, "02  Información dispersa",
         "Los datos del paciente y sus atenciones quedan sin una trazabilidad central.", GOLD, True)
add_card(slide, 0.65, 4.66, 6.35, 1.10, "03  Poco control operativo",
         "Es difícil supervisar agendas, servicios, ingresos y rendimiento de la clínica.", RED, True)
add_footer(slide, 3, True)


# 4. Solución y objetivos
slide = prs.slides.add_slide(blank)
add_header(slide, "Solución propuesta", "Una plataforma central para pacientes, odontólogos y administración", 4)
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.66, 1.95, 4.0, 4.55, NAVY, NAVY, True)
add_text(slide, "OBJETIVO GENERAL", 1.02, 2.28, 3.25, 0.28, 11, GOLD_LIGHT, True)
add_text(slide, "Desarrollar un sistema web integrado que gestione de forma segura y eficiente los procesos odontológicos.",
         1.02, 2.80, 3.25, 1.48, 21, WHITE, True, line_spacing=1.05)
add_text(slide, "Resultado esperado", 1.02, 4.72, 3.25, 0.30, 11, CYAN, True)
add_text(slide, "Menos errores, mejor trazabilidad y una experiencia clara para cada rol.",
         1.02, 5.12, 3.15, 0.88, 13, "C6D8EB")
objectives = [
    ("API RESTful", "Servicios organizados y documentados."),
    ("Persistencia", "Datos reales mediante JPA y MariaDB."),
    ("Seguridad", "Autenticación JWT y permisos por rol."),
    ("Integración", "Frontend conectado a la API del backend."),
]
for i, (title, body) in enumerate(objectives):
    x = 5.08 + (i % 2) * 3.76
    y = 2.02 + (i // 2) * 2.18
    add_card(slide, x, y, 3.40, 1.78, title, body, CYAN if i % 2 == 0 else GOLD)
add_footer(slide, 4)


# 5. Arquitectura
slide = prs.slides.add_slide(blank)
add_header(slide, "Arquitectura del sistema", "Separación por capas para mantener responsabilidades claras", 5)
layers = [
    ("FRONTEND", "HTML · Tailwind CSS · JavaScript", CYAN),
    ("API REST", "Controladores + DTO + validaciones", GOLD),
    ("NEGOCIO", "Servicios + reglas + seguridad JWT", "4B7BEC"),
    ("PERSISTENCIA", "Spring Data JPA · Hibernate", "9B59B6"),
    ("DATOS", "MariaDB · smartdent_db", GREEN),
]
for i, (title, body, accent) in enumerate(layers):
    x = 0.85 + i * 2.47
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.56, 2.02, 2.24, WHITE, LINE, True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, 2.56, 2.02, 0.12, accent)
    add_text(slide, f"0{i+1}", x + 0.18, 2.86, 0.35, 0.35, 11, accent, True)
    add_text(slide, title, x + 0.18, 3.30, 1.66, 0.35, 14, NAVY, True,
             align=PP_ALIGN.CENTER)
    add_text(slide, body, x + 0.18, 3.82, 1.66, 0.62, 10.5, MUTED,
             align=PP_ALIGN.CENTER)
    if i < len(layers) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x + 2.06), Inches(3.41),
                                       Inches(0.36), Inches(0.46))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = rgb(GOLD)
        arrow.line.color.rgb = rgb(GOLD)
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 2.13, 5.48, 9.05, 0.72, "EAFBFD", "BCEFF3", True)
add_text(slide, "Flujo principal: navegador → HTTP/JSON → Controller → Service → Repository → MariaDB",
         2.35, 5.71, 8.62, 0.25, 12, NAVY, True, align=PP_ALIGN.CENTER)
add_footer(slide, 5)


# 6. Roles
slide = prs.slides.add_slide(blank)
add_header(slide, "Tres roles, una sola plataforma", "Cada usuario accede únicamente a las funciones que necesita", 6)
roles = [
    ("PACIENTE", "P", CYAN, ["Reserva y reprograma citas", "Consulta historial clínico", "Configura recordatorios"]),
    ("ODONTÓLOGO", "O", GOLD, ["Gestiona su agenda", "Registra diagnóstico y tratamiento", "Bloquea horarios"]),
    ("ADMINISTRADOR", "A", "4B7BEC", ["Supervisa la agenda global", "Gestiona profesionales y tarifas", "Analiza reportes"]),
]
for i, (title, letter, accent, bullets) in enumerate(roles):
    x = 0.67 + i * 4.2
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.0, 3.78, 4.40, WHITE, LINE, True)
    add_shape(slide, MSO_SHAPE.OVAL, x + 1.42, 2.36, 0.92, 0.92, accent)
    add_text(slide, letter, x + 1.42, 2.52, 0.92, 0.40, 22, WHITE, True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, title, x + 0.30, 3.53, 3.18, 0.42, 16, NAVY, True,
             align=PP_ALIGN.CENTER)
    add_bullet_list(slide, bullets, x + 0.38, 4.18, 3.02, 1.62, 11.5, MUTED, accent, 0.58)
add_footer(slide, 6)


# 7. Flujo de cita
slide = prs.slides.add_slide(blank)
add_header(slide, "Flujo principal: reservar una cita", "Proceso completo y trazable desde el catálogo hasta la atención", 7)
steps = [
    ("1", "Servicio", "El paciente elige el tratamiento."),
    ("2", "Profesional", "Selecciona un odontólogo disponible."),
    ("3", "Fecha y hora", "La API valida cruces y bloqueos."),
    ("4", "Confirmación", "La cita queda en estado pendiente."),
    ("5", "Atención", "El odontólogo registra la historia clínica."),
]
for i, (num, title, body) in enumerate(steps):
    x = 0.68 + i * 2.50
    add_shape(slide, MSO_SHAPE.OVAL, x + 0.65, 2.14, 0.82, 0.82, NAVY)
    add_text(slide, num, x + 0.65, 2.31, 0.82, 0.32, 17, WHITE, True,
             align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_shape(slide, MSO_SHAPE.RECTANGLE, x + 1.46, 2.49, 1.69, 0.10, GOLD)
    add_text(slide, title, x, 3.28, 2.16, 0.35, 14, NAVY, True,
             align=PP_ALIGN.CENTER)
    add_text(slide, body, x, 3.80, 2.16, 0.92, 10.5, MUTED,
             align=PP_ALIGN.CENTER)
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.16, 5.35, 11.0, 0.86, "FFF9DF", "E9D578", True)
add_text(slide, "Regla clave: un horario reservado o bloqueado desaparece de la disponibilidad y no puede duplicarse.",
         1.46, 5.62, 10.40, 0.28, 12.2, NAVY, True, align=PP_ALIGN.CENTER)
add_footer(slide, 7)


# 8. Persistencia
slide = prs.slides.add_slide(blank)
add_header(slide, "Persistencia y modelo de datos", "Los datos de negocio se almacenan en MariaDB, no en el navegador", 8, True)
entities = [
    ("USUARIOS", "roles · configuración"),
    ("ODONTÓLOGOS", "servicios asignados"),
    ("CITAS", "estado · fecha · precio"),
    ("HISTORIAS", "diagnóstico · tratamiento"),
    ("SERVICIOS", "tarifa · costo · duración"),
    ("REPORTES", "costos fijos · indicadores"),
]
for i, (title, body) in enumerate(entities):
    col, row = i % 3, i // 3
    x, y = 0.72 + col * 4.14, 2.04 + row * 1.70
    add_card(slide, x, y, 3.72, 1.36, title, body, CYAN if row == 0 else GOLD, True)
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.42, 5.66, 10.5, 0.74, NAVY_2, "244263", True)
add_rich_text(slide, [("Navegador: ", CYAN, True), ("solo JWT y navegación temporal   ", WHITE, False),
                      ("Base de datos: ", GOLD_LIGHT, True), ("usuarios, citas, historias, servicios y reportes", WHITE, False)],
              1.73, 5.89, 9.9, 0.28, 11.5, align=PP_ALIGN.CENTER)
add_footer(slide, 8, True)


# 9. API y seguridad
slide = prs.slides.add_slide(blank)
add_header(slide, "API REST, documentación y seguridad", "Endpoints organizados por contexto y protegidos según el rol", 9)
groups = [
    ("PÚBLICO", "/api/auth · /api/servicios\n/api/odontologos · /api/contacto", CYAN),
    ("PACIENTE", "/api/pacientes/citas\n/historias-clinicas · configuración", GOLD),
    ("ODONTÓLOGO", "/api/odontologos/mi-agenda\n/citas · historias · bloqueos", "4B7BEC"),
    ("ADMIN", "/api/admin/citas · usuarios\nservicios · reportes · mensajes", "9B59B6"),
]
for i, (title, body, accent) in enumerate(groups):
    x = 0.68 + i * 3.11
    add_card(slide, x, 1.98, 2.76, 1.75, title, body, accent)
add_text(slide, "Autenticación", 0.72, 4.22, 2.1, 0.33, 15, NAVY, True)
security_steps = ["Correo + contraseña", "BCrypt verifica", "JWT firmado", "Permiso por rol"]
for i, label in enumerate(security_steps):
    x = 0.73 + i * 3.04
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 4.83, 2.42, 0.78, NAVY if i != 2 else GOLD,
              NAVY if i != 2 else GOLD, True)
    add_text(slide, label, x + 0.12, 5.06, 2.18, 0.28, 11, WHITE, True,
             align=PP_ALIGN.CENTER)
    if i < 3:
        add_text(slide, "→", x + 2.52, 5.02, 0.36, 0.32, 18, GOLD, True,
                 align=PP_ALIGN.CENTER)
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 3.78, 6.10, 5.78, 0.54, "EAFBFD", "BCEFF3", True)
add_text(slide, "Swagger UI: http://localhost:8080/swagger-ui.html", 3.98, 6.25, 5.38, 0.23,
         10.8, NAVY, True, align=PP_ALIGN.CENTER)
add_footer(slide, 9)


# 10. TDD y resultados
slide = prs.slides.add_slide(blank)
add_header(slide, "TDD y estrategia de pruebas", "Retroalimentación continua: rojo, verde y refactorización", 10, True)
cycles = [
    ("ROJO", "Escribir una prueba que falle", RED),
    ("VERDE", "Implementar lo mínimo para aprobar", GREEN),
    ("REFACTOR", "Mejorar sin cambiar el resultado", "4B7BEC"),
]
for i, (title, body, accent) in enumerate(cycles):
    x = 0.75 + i * 3.10
    add_shape(slide, MSO_SHAPE.OVAL, x, 2.16, 2.25, 2.25, accent)
    add_text(slide, title, x + 0.25, 2.72, 1.75, 0.36, 16, WHITE, True,
             align=PP_ALIGN.CENTER)
    add_text(slide, body, x + 0.27, 3.20, 1.71, 0.62, 10.5, WHITE,
             align=PP_ALIGN.CENTER)
    if i < 2:
        add_text(slide, "→", x + 2.38, 3.03, 0.46, 0.35, 22, GOLD_LIGHT, True,
                 align=PP_ALIGN.CENTER)
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 10.35, 1.95, 2.20, 3.00, WHITE, WHITE, True)
add_text(slide, "34", 10.46, 2.31, 1.98, 0.78, 40, NAVY, True, align=PP_ALIGN.CENTER)
add_text(slide, "PRUEBAS\nAPROBADAS", 10.58, 3.18, 1.74, 0.72, 13, GREEN, True,
         align=PP_ALIGN.CENTER)
add_text(slide, "0 fallos · 0 errores", 10.53, 4.25, 1.84, 0.28, 10.5, MUTED, True,
         align=PP_ALIGN.CENTER)
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.76, 5.36, 11.78, 0.96, NAVY_2, "244263", True)
add_text(slide, "Cobertura principal: autenticación · roles · citas · disponibilidad · historias clínicas · reportes · OpenAPI",
         1.08, 5.68, 11.14, 0.30, 11.7, "D7E6F7", True, align=PP_ALIGN.CENTER)
add_footer(slide, 10, True)


# 11. Código JWT
slide = prs.slides.add_slide(blank)
add_header(slide, "Prueba clave 01 · Seguridad JWT", "Verifica que el login genere un token válido con el rol correcto", 11)
code = """@Test
void debeIniciarSesionYGenerarUnJwtValido() {
    Credenciales c = registrarPaciente();

    LoginResponse response = loginService.iniciarSesion(
        new LoginRequest(c.email(), c.password()));
    Jwt jwt = jwtDecoder.decode(response.token());

    assertThat(response.tokenType()).isEqualTo(\"Bearer\");
    assertThat(response.expiresIn()).isPositive();
    assertThat(jwt.getSubject()).isEqualTo(c.email());
    assertThat(jwt.getClaimAsStringList(\"rol\"))
        .containsExactly(\"PACIENTE\");
}"""
add_code_block(slide, code, 0.62, 1.95, 8.45, 4.82)
add_card(slide, 9.38, 1.95, 3.30, 1.22, "QUÉ COMPRUEBA", "El token pertenece al usuario autenticado.", CYAN)
add_card(slide, 9.38, 3.42, 3.30, 1.22, "REGLA DE SEGURIDAD", "El rol PACIENTE viaja como claim del JWT.", GOLD)
add_card(slide, 9.38, 4.89, 3.30, 1.22, "RESULTADO", "La sesión puede autorizar endpoints protegidos.", GREEN)
add_text(slide, "AuthSecurityIntegrationTests.java", 9.52, 6.42, 3.03, 0.23, 9, MUTED,
         align=PP_ALIGN.CENTER)
add_footer(slide, 11)


# 12. Código citas
slide = prs.slides.add_slide(blank)
add_header(slide, "Prueba clave 02 · Reserva sin conflictos", "Comprueba persistencia, visibilidad por rol y prevención de cruces", 12)
code = """@Test
void debeReservarYMostrarLaCitaEnLosTresPaneles() {
    var cita = citaService.reservar(paciente,
        new CrearCitaRequest(odontologo.getId(),
            servicio.getId(), fecha, LocalTime.of(10, 0),
            \"Evaluación preventiva\", \"987654321\"));

    assertThat(cita.estado()).isEqualTo(PENDIENTE);
    assertThat(citaService.listarDelPaciente(paciente))
        .extracting(\"id\").contains(cita.id());
    assertThat(citaService.listarDelOdontologo(correoDoctor))
        .extracting(\"id\").contains(cita.id());
    assertThat(citaService.listarTodas())
        .extracting(\"id\").contains(cita.id());
}"""
add_code_block(slide, code, 0.62, 1.95, 8.45, 4.82)
add_card(slide, 9.38, 1.95, 3.30, 1.22, "PACIENTE", "La nueva cita aparece en su panel.", CYAN)
add_card(slide, 9.38, 3.42, 3.30, 1.22, "ODONTÓLOGO", "Solo el profesional asignado puede verla.", GOLD)
add_card(slide, 9.38, 4.89, 3.30, 1.22, "ADMIN", "La reserva aparece en la agenda global.", "4B7BEC")
add_text(slide, "CitaServiceIntegrationTests.java · fragmento real simplificado", 9.40, 6.38, 3.25, 0.34,
         8.7, MUTED, align=PP_ALIGN.CENTER)
add_footer(slide, 12)


# 13. Código historia clínica
slide = prs.slides.add_slide(blank)
add_header(slide, "Prueba clave 03 · Historia clínica", "La atención actualiza el expediente y marca la cita como atendida", 13)
code = """@Test
void debeGuardarLaHistoriaYMarcarLaCitaAtendida() {
    citaService.cambiarEstadoPorOdontologo(
        correoDoctor, cita.id(), CONFIRMADA);

    var historia = historiaClinicaService.guardar(
        correoDoctor, paciente,
        new GuardarHistoriaClinicaRequest(
            cita.id(), TRATAMIENTO, \"Ninguna\",
            \"Pulpitis irreversible\", \"Endodoncia iniciada\",
            \"Tomar el medicamento indicado\", control, null));

    assertThat(historia.ultimaCitaId()).isEqualTo(cita.id());
    assertThat(citaService.listarDelPaciente(paciente)
        .getFirst().estado()).isEqualTo(ATENDIDA);
}"""
add_code_block(slide, code, 0.62, 1.95, 8.45, 4.82)
add_card(slide, 9.38, 1.95, 3.30, 1.22, "CONDICIÓN", "La cita debe estar confirmada antes de atenderse.", GOLD)
add_card(slide, 9.38, 3.42, 3.30, 1.22, "TRAZABILIDAD", "Diagnóstico y tratamiento quedan vinculados.", CYAN)
add_card(slide, 9.38, 4.89, 3.30, 1.22, "ESTADO", "Al guardar, la cita cambia a ATENDIDA.", GREEN)
add_text(slide, "HistoriaClinicaIntegrationTests.java · fragmento real simplificado", 9.40, 6.38, 3.25, 0.34,
         8.5, MUTED, align=PP_ALIGN.CENTER)
add_footer(slide, 13)


# 14. Cierre
slide = prs.slides.add_slide(blank)
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 7.5, NAVY)
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.08, CYAN)
add_logo(slide, 0.70, 0.48, 2.42)
add_text(slide, "Resultados del avance", 0.72, 1.53, 6.3, 0.62, 31, WHITE, True)
add_text(slide, "Una base funcional, segura y preparada para evolucionar.", 0.74, 2.24, 6.2, 0.48,
         16, "BED0E4")
results = [
    "Frontend conectado a una API REST real.",
    "Persistencia central en MariaDB.",
    "Paneles separados para tres roles.",
    "JWT, BCrypt y permisos por endpoint.",
    "34 pruebas automatizadas aprobadas.",
]
add_bullet_list(slide, results, 0.78, 3.02, 5.70, 2.80, 14, WHITE, GOLD, 0.54)
add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.32, 1.15, 5.20, 5.28, WHITE, WHITE, True)
add_text(slide, "PRÓXIMOS PASOS", 7.78, 1.65, 4.28, 0.32, 12, GOLD, True)
next_steps = [
    "Migrar el frontend a Angular.",
    "Incorporar validaciones y componentes reutilizables.",
    "Reforzar seguridad para producción con HTTPS.",
    "Desplegar frontend, API y base de datos en la nube.",
]
add_bullet_list(slide, next_steps, 7.80, 2.30, 4.10, 2.48, 13, NAVY, CYAN, 0.62)
add_shape(slide, MSO_SHAPE.RECTANGLE, 7.80, 5.22, 3.80, 0.05, GOLD)
add_text(slide, "Gracias", 7.80, 5.52, 4.10, 0.46, 23, NAVY, True)
add_text(slide, "¿Preguntas?", 7.80, 6.00, 4.10, 0.30, 12, MUTED)
add_text(slide, "14", 12.24, 7.13, 0.50, 0.20, 8, "8DA2BD", True,
         align=PP_ALIGN.RIGHT)


OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"Presentación creada: {OUT}")
print(f"Diapositivas: {len(prs.slides)}")
